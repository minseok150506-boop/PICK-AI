from functools import wraps

def thinking_process(text: str):
    steps = []
    steps.append(f"요청 이해: {text}")

    if "ppt" in text.lower():
        steps.append("작업 분석: PPT 생성 요청")
        steps.append("처리 계획: 주제 정리 → 슬라이드 구성 → 생성")
    elif any(op in text for op in "+-*/"):
        steps.append("작업 분석: 계산 요청")
        steps.append("처리 계획: 수식 처리 후 결과 반환")
    else:
        steps.append("작업 분석: 일반 질문")
        steps.append("처리 계획: 정보 정리 후 답변")

    return "\n".join(steps)


import io
import os
import ast
import operator as op
import re

import re
import sqlite3
from datetime import datetime
from pathlib import Path
from pick_engine import PickEngine
from pick_llm import PickLLMRouter
from pick_typo import PickTypoEngine
from pick_polish import final_review

from flask import Flask, jsonify, redirect, render_template, request, send_file, send_from_directory, session, url_for
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename

try:
    from openai import OpenAI
except Exception:
    OpenAI = None

try:
    from PIL import Image, ImageStat
except Exception:
    Image = None
    ImageStat = None

try:
    from pptx import Presentation
    from pptx.util import Pt
except Exception:
    Presentation = None
    Pt = None

APP_DIR = Path(__file__).resolve().parent
STORAGE = APP_DIR / "storage"
UPLOADS = STORAGE / "uploads"
GENERATED = STORAGE / "generated"
DB_PATH = STORAGE / "pick.db"

UPLOADS.mkdir(parents=True, exist_ok=True)
GENERATED.mkdir(parents=True, exist_ok=True)

SITE_RECOMMENDATIONS = {
    "이미지": [
        {"name": "Unsplash", "url": "https://unsplash.com", "desc": "무료 고품질 사진 사이트"},
        {"name": "Pexels", "url": "https://www.pexels.com", "desc": "무료 사진과 영상 사이트"},
        {"name": "Pixabay", "url": "https://pixabay.com", "desc": "무료 이미지, 벡터, 영상 제공"}
    ],
    "ppt": [
        {"name": "Slidesgo", "url": "https://slidesgo.com", "desc": "PPT 템플릿 사이트"},
        {"name": "Canva", "url": "https://www.canva.com", "desc": "디자인과 발표자료 제작"},
        {"name": "Beautiful.ai", "url": "https://www.beautiful.ai", "desc": "AI 기반 발표자료 제작"}
    ],
    "코딩": [
        {"name": "W3Schools", "url": "https://www.w3schools.com", "desc": "기초 코딩 학습"},
        {"name": "MDN", "url": "https://developer.mozilla.org", "desc": "웹 개발 공식 문서"},
        {"name": "GitHub", "url": "https://github.com", "desc": "코드 저장소와 오픈소스"}
    ],
    "게임에셋": [
        {"name": "itch.io", "url": "https://itch.io/game-assets", "desc": "게임 에셋 자료"},
        {"name": "Kenney", "url": "https://kenney.nl/assets", "desc": "무료 게임 에셋"},
        {"name": "OpenGameArt", "url": "https://opengameart.org", "desc": "오픈 게임 아트 자료"}
    ]
}

def recommend_sites(text: str):
    lowered = text.lower()
    if ("사이트" in text or "추천" in text or "찾아" in text or "알려" in text):
        if "이미지" in text or "사진" in text:
            return SITE_RECOMMENDATIONS["이미지"]
        if "ppt" in lowered or "피피티" in text or "발표" in text or "템플릿" in text:
            return SITE_RECOMMENDATIONS["ppt"]
        if "코딩" in text or "프로그래밍" in text or "개발" in text:
            return SITE_RECOMMENDATIONS["코딩"]
        if "게임" in text and ("에셋" in text or "asset" in lowered):
            return SITE_RECOMMENDATIONS["게임에셋"]
    return []


app = Flask(__name__)


# ===== PICK SERVICE CLEAN AUTH HELPERS =====
DB_PATH = "data/pick_service.db"

def db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def now_iso():
    return datetime.now().isoformat(timespec="seconds")

def init_service_db():
    Path("data").mkdir(exist_ok=True)
    conn = db()
    cur = conn.cursor()
    cur.execute("""
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        password_hash TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)
    cur.execute("""
    CREATE TABLE IF NOT EXISTS conversations (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        title TEXT NOT NULL,
        created_at TEXT NOT NULL,
        updated_at TEXT NOT NULL
    )
    """)
    cur.execute("""
    CREATE TABLE IF NOT EXISTS messages (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        conversation_id INTEGER NOT NULL,
        role TEXT NOT NULL,
        content TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)
    cur.execute("""
    CREATE TABLE IF NOT EXISTS service_logs (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        level TEXT NOT NULL,
        message TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)
    conn.commit()
    conn.close()

def service_log(level, message):
    try:
        conn = db()
        conn.execute(
            "INSERT INTO service_logs(level, message, created_at) VALUES (?, ?, ?)",
            (level, str(message), now_iso())
        )
        conn.commit()
        conn.close()
    except Exception:
        pass

def login_required(fn):
    @wraps(fn)
    def wrapper(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        return fn(*args, **kwargs)
    return wrapper

def has_korean(text):
    return any("가" <= ch <= "힣" for ch in str(text or ""))

def valid_account_text(text, allow_symbols=False):
    import re
    if allow_symbols:
        return re.fullmatch(r"[A-Za-z0-9!@#$%^&*()_\-+=.?]{4,64}", text or "") is not None
    return re.fullmatch(r"[A-Za-z0-9_\-]{2,32}", text or "") is not None

init_service_db()



# ===== PICK SERVICE AUTH + DB LAYER =====
import sqlite3
from functools import wraps
from datetime import datetime
from werkzeug.security import generate_password_hash, check_password_hash

DB_PATH = "data/pick_service.db"

def db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_service_db():
    Path("data").mkdir(exist_ok=True)
    conn = db()
    cur = conn.cursor()

    cur.execute("""
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username TEXT UNIQUE NOT NULL,
        password_hash TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)

    cur.execute("""
    CREATE TABLE IF NOT EXISTS conversations (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        title TEXT NOT NULL,
        created_at TEXT NOT NULL,
        updated_at TEXT NOT NULL,
        FOREIGN KEY(user_id) REFERENCES users(id)
    )
    """)

    cur.execute("""
    CREATE TABLE IF NOT EXISTS messages (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        conversation_id INTEGER NOT NULL,
        role TEXT NOT NULL,
        content TEXT NOT NULL,
        created_at TEXT NOT NULL,
        FOREIGN KEY(conversation_id) REFERENCES conversations(id)
    )
    """)

    cur.execute("""
    CREATE TABLE IF NOT EXISTS service_logs (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        level TEXT NOT NULL,
        message TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)

    conn.commit()
    conn.close()

def now_iso():
    return datetime.now().isoformat(timespec="seconds")

def service_log(level, message):
    try:
        conn = db()
        conn.execute(
            "INSERT INTO service_logs(level, message, created_at) VALUES (?, ?, ?)",
            (level, str(message), now_iso())
        )
        conn.commit()
        conn.close()
    except Exception:
        pass

def login_required(fn):
    @wraps(fn)
    def wrapper(*args, **kwargs):
        if not session.get("user_id"):
            return redirect(url_for("login"))
        return fn(*args, **kwargs)
    return wrapper

def current_user():
    uid = session.get("user_id")
    if not uid:
        return None
    conn = db()
    user = conn.execute("SELECT id, username, created_at FROM users WHERE id=?", (uid,)).fetchone()
    conn.close()
    return user

init_service_db()
pick_engine = PickEngine()
pick_llm = PickLLMRouter()
pick_typo = PickTypoEngine()
app.secret_key = os.environ.get("PICK_SECRET_KEY", "pick_secret_change_me")
ADMIN_USERNAME = os.environ.get("PICK_ADMIN_USERNAME", "minseok")
ADMIN_PASSWORD = os.environ.get("PICK_ADMIN_PASSWORD", "kms0506a!")

user_states = {}

def get_state(user_id: int):
    if user_id not in user_states:
        user_states[user_id] = {"task": None, "topic": None, "options": {}}
    return user_states[user_id]

def db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = db()
    cur = conn.cursor()
    cur.execute("""
    CREATE TABLE IF NOT EXISTS users (
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      username TEXT UNIQUE NOT NULL,
      password_hash TEXT NOT NULL,
      created_at TEXT NOT NULL
    )""")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS chats (
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      user_id INTEGER NOT NULL,
      title TEXT NOT NULL DEFAULT '새 채팅',
      created_at TEXT NOT NULL
    )""")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS messages (
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      chat_id INTEGER NOT NULL,
      role TEXT NOT NULL,
      content TEXT NOT NULL,
      created_at TEXT NOT NULL
    )""")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS projects (
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      user_id INTEGER NOT NULL,
      title TEXT NOT NULL,
      description TEXT NOT NULL,
      created_at TEXT NOT NULL
    )""")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS requests (
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      user_id INTEGER NOT NULL,
      content TEXT NOT NULL,
      created_at TEXT NOT NULL
    )""")

    # 관리자 계정은 새 데이터베이스에서도 회원가입 없이 바로 로그인할 수 있도록 보장합니다.
    admin = cur.execute("SELECT id FROM users WHERE username=?", (ADMIN_USERNAME,)).fetchone()
    admin_hash = generate_password_hash(ADMIN_PASSWORD)
    if admin:
        cur.execute("UPDATE users SET password_hash=? WHERE id=?", (admin_hash, admin["id"]))
    else:
        cur.execute(
            "INSERT INTO users(username, password_hash, created_at) VALUES(?,?,?)",
            (ADMIN_USERNAME, admin_hash, now())
        )

    conn.commit()
    conn.close()

def current_user_id():
    return session.get("user_id")

def logged_in():
    return current_user_id() is not None

def require_login_json():
    if not logged_in():
        return jsonify({"ok": False, "error": "로그인이 필요합니다."}), 401
    return None

def now():
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")

def get_user_chats(user_id):
    conn = db()
    rows = conn.execute("SELECT * FROM chats WHERE user_id=? ORDER BY id DESC", (user_id,)).fetchall()
    conn.close()
    return [dict(r) for r in rows]

def get_messages(chat_id):
    conn = db()
    rows = conn.execute("SELECT * FROM messages WHERE chat_id=? ORDER BY id ASC", (chat_id,)).fetchall()
    conn.close()
    return [dict(r) for r in rows]

def create_chat(user_id, with_greeting=True):
    conn = db()
    cur = conn.cursor()
    cur.execute("INSERT INTO chats(user_id, title, created_at) VALUES(?,?,?)", (user_id, "새 채팅", now()))
    chat_id = cur.lastrowid
    if with_greeting:
        cur.execute("INSERT INTO messages(chat_id, role, content, created_at) VALUES(?,?,?,?)",
                    (chat_id, "bot", "안녕하세요 저는 PICK 입니다 무엇을 도와드릴까요?", now()))
    conn.commit()
    conn.close()
    return chat_id

def update_chat_title(chat_id):
    conn = db()
    row = conn.execute("SELECT content FROM messages WHERE chat_id=? AND role='user' ORDER BY id ASC LIMIT 1", (chat_id,)).fetchone()
    title = (row["content"][:20] if row else "새 채팅")
    conn.execute("UPDATE chats SET title=? WHERE id=?", (title, chat_id))
    conn.commit()
    conn.close()

def is_ppt_request(text):
    return bool(re.search(r"(ppt|피피티|프레젠테이션|발표자료)", text, re.I) and re.search(r"(만들|작성|생성|준비)", text, re.I))

def topic_of(text):
    cleaned = re.sub(r"ppt|피피티|프레젠테이션|발표자료|만들어줘|만들기|만들어|작성해줘|생성해줘|준비해줘", " ", text, flags=re.I)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if "호랑이" in cleaned:
        return "호랑이"
    return cleaned or "주제"

def classify_topic(topic):
    t = topic.strip()
    if any(k in t for k in ["호랑이", "동물", "사자", "고양이", "개", "생물"]):
        return "animal"
    if any(k in t for k in ["자동차", "차", "전기차", "엔진"]):
        return "vehicle"
    if any(k in t for k in ["게임", "기획", "게임기획", "RPG", "액션"]):
        return "game"
    if any(k in t for k in ["우주", "행성", "은하", "별"]):
        return "space"
    return "general"

def deck_data(topic, options=None):
    options = options or {}
    category = classify_topic(topic)
    slides_override = options.get("slides")
    title = options.get("title") or topic

    if category == "animal":
        slides = [
            ("정의", [f"{topic}의 기본 정의를 설명합니다.", "생물학적 분류와 기본 개념을 정리합니다.", "대표적인 특징을 간단히 소개합니다."]),
            ("핵심 특징", ["외형적 특징을 설명합니다.", "행동이나 습성을 요약합니다.", "다른 동물과 구분되는 점을 정리합니다."]),
            ("기본 구성 요소", ["서식지", "먹이", "생활 방식"]),
            ("중요성", ["생태계에서의 역할을 설명합니다.", "보호가 필요한 이유를 정리합니다.", "인간과의 관계를 간단히 설명합니다."]),
            ("정리", [f"{topic}의 핵심 내용을 다시 요약합니다.", "기억할 포인트를 정리합니다.", "추가 탐구 방향을 제안합니다."]),
        ]
    elif category == "vehicle":
        slides = [
            ("정의", [f"{topic}의 개념과 목적을 설명합니다.", "기본 역할을 정리합니다.", "대표 예시를 간단히 제시합니다."]),
            ("핵심 특징", ["속도/동력/이동 방식 등 핵심 특징을 설명합니다.", "용도별 차이를 간단히 정리합니다.", "대표 장점을 요약합니다."]),
            ("기본 구성 요소", ["차체", "동력 장치", "제어 장치"]),
            ("활용", ["어디에 쓰이는지 설명합니다.", "실생활/산업적 중요성을 정리합니다.", "대표 활용 예를 제시합니다."]),
            ("정리", [f"{topic}의 주요 내용을 다시 정리합니다.", "핵심 요소를 한 번 더 요약합니다.", "추가 확장 포인트를 제안합니다."]),
        ]
    elif category == "game":
        slides = [
            ("개요", [f"{topic}의 전체 개념을 정리합니다.", "목표와 방향을 설명합니다.", "핵심 장르나 성격을 요약합니다."]),
            ("핵심 특징", ["장르", "핵심 재미", "주요 차별점을 정리합니다."]),
            ("기본 구성 요소", ["플레이 방식", "성장 구조", "콘텐츠 구조"]),
            ("활용/중요성", ["대상 유저", "시장성", "운영 방향을 설명합니다."]),
            ("정리", [f"{topic} 기획의 핵심을 다시 정리합니다.", "중요 포인트를 요약합니다.", "다음 단계 제안을 넣습니다."]),
        ]
    elif category == "space":
        slides = [
            ("정의", [f"{topic}의 기본 개념을 설명합니다.", "관련 과학 개념을 정리합니다.", "주요 범위를 소개합니다."]),
            ("핵심 특징", ["구성 요소를 설명합니다.", "대표 현상을 정리합니다.", "관찰 포인트를 소개합니다."]),
            ("기본 구성 요소", ["구조", "원리", "대표 사례"]),
            ("중요성", ["과학적 의미를 설명합니다.", "연구 가치와 활용을 정리합니다.", "현대 기술과의 연결점을 설명합니다."]),
            ("정리", [f"{topic}의 핵심 내용을 요약합니다.", "기억할 포인트를 정리합니다.", "확장 학습 방향을 제안합니다."]),
        ]
    else:
        slides = [
            ("정의", [f"{topic}의 정의와 핵심 개념을 정리합니다.", "기본 배경을 설명합니다.", "전체 범위를 간단히 소개합니다."]),
            ("핵심 특징", [f"{topic}의 주요 특징을 설명합니다.", "대표 성질이나 성격을 정리합니다.", "핵심 포인트를 요약합니다."]),
            ("기본 구성 요소", [f"{topic}를 이해하는 데 필요한 기본 요소를 설명합니다.", "중요 구성 요소를 나열합니다.", "기초 이해에 필요한 개념을 넣습니다."]),
            ("활용/중요성", [f"{topic}가 왜 중요한지 설명합니다.", "어디에 활용되는지 정리합니다.", "실제 의미를 설명합니다."]),
            ("정리", [f"{topic}의 핵심 내용을 다시 요약합니다.", "기억할 포인트를 정리합니다.", "추가 확장 방향을 제안합니다."]),
        ]

    if "호랑이" in topic:
        slides = [
            ("호랑이란?", ["포유류 식육목 고양이과에 속하는 대형 맹수입니다.", "검은 줄무늬가 있는 주황빛 털이 대표 특징입니다.", "주로 단독 생활을 합니다."]),
            ("기본 특징", ["학명: Panthera tigris", "서식지: 아시아의 숲, 초원, 습지", "먹이: 사슴, 멧돼지 등"]),
            ("몸의 구조", ["강한 턱과 날카로운 이빨을 가졌습니다.", "근육질 몸과 강한 다리로 사냥합니다.", "줄무늬는 위장에 도움이 됩니다."]),
            ("생활 방식", ["주로 혼자 생활합니다.", "야간에 활발히 움직이는 편입니다.", "수영을 잘합니다."]),
            ("보호 필요성", ["서식지 파괴와 밀렵으로 개체 수가 감소했습니다.", "생태계 균형에 중요한 동물입니다.", "보호 정책과 국제 협력이 필요합니다."]),
        ]

    if isinstance(slides_override, int) and slides_override > 0:
        if slides_override < len(slides):
            slides = slides[:slides_override]
        elif slides_override > len(slides):
            base_slides = list(slides)
            i = 1
            while len(base_slides) < slides_override:
                base_slides.append((f"추가 정리 {i}", [f"{topic}의 추가 설명 {i}입니다.", "핵심 내용을 보강합니다.", "요약 정리를 이어갑니다."]))
                i += 1
            slides = base_slides

    return {"title": title, "slides": slides}

def create_ppt(topic, options=None):
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다.")
    data = deck_data(topic, options)
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = data["title"]
    title_slide.placeholders[1].text = "PICK 자동 생성 발표자료"

    for title, bullets in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(20)

    safe_name = re.sub(r"[^\w가-힣\- ]", "_", data["title"])
    path = GENERATED / f"{safe_name}_PICK.pptx"
    prs.save(path)
    return path.name

def analyze_image(path):
    if Image is None:
        return {"error": "Pillow가 설치되어 있지 않습니다."}
    img = Image.open(path).convert("RGB")
    stat = ImageStat.Stat(img)
    mean = tuple(int(v) for v in stat.mean[:3])
    brightness = round(sum(mean) / 3, 1)
    return {
        "filename": path.name,
        "size_bytes": path.stat().st_size,
        "width": img.width,
        "height": img.height,
        "ratio": round(img.width / img.height, 2) if img.height else None,
        "mean_rgb": mean,
        "brightness": brightness,
        "note": "업로드한 아무 이미지를 구조 정보 중심으로 분석합니다."
    }

def analyze_file(path):
    st = path.stat()
    return {
        "filename": path.name,
        "extension": path.suffix,
        "size_bytes": st.st_size,
        "modified": datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
        "path": str(path)
    }


def should_use_ai(text: str):
    text = (text or "").strip()
    lowered = text.lower()

    # 도구 실행형 명령은 로컬 우선
    if "사이트 추천" in text:
        return False
    if try_calculate(text):
        return False
    if any(k in lowered for k in ["ppt", "exe"]) or any(k in text for k in ["피피티", "이미지 분석", "동영상 분석", "파일 분석", "프로젝트", "요청"]):
        return False

    # 일반 대화는 거의 전부 AI 사용
    return len(text) >= 1
    if any(k in lowered for k in ["ppt", "exe"]) or any(k in text for k in ["피피티", "이미지 분석", "동영상 분석", "파일 분석", "프로젝트", "요청"]):
        return False
    if len(text) >= 24:
        return True
    return False

def local_chat_reply(text: str, username: str):
    text = (text or "").strip()
    lowered = text.lower()
    if not text:
        return "말씀을 입력해 주세요."
    if "안녕" in text:
        return "안녕하세요. 무엇을 도와드릴까요?"
    if "이름" in text:
        return f"현재 로그인한 아이디는 {username} 입니다."
    if "날씨" in text:
        return "날씨 기능은 아직 실시간 연결이 없습니다. 지역을 말씀해 주시면 날씨 기능 방향을 잡을 수 있습니다."
    if "이미지" in text and any(k in text for k in ["열어", "분석", "화면", "기능"]):
        return "이미지 분석 화면을 열겠습니다."
    if (("동영상" in text or "영상" in text) and any(k in text for k in ["열어", "분석", "화면", "기능"])):
        return "동영상 분석 화면을 열겠습니다."
    if "파일" in text and any(k in text for k in ["열어", "분석", "화면", "기능"]):
        return "파일 분석 화면을 열겠습니다."
    if "ppt" in lowered or "피피티" in text:
        return "PPT 작업을 시작하겠습니다. 주제와 장수를 말씀해 주세요."
    return "요청을 이해했습니다. 조금 더 구체적으로 말씀해 주세요."


_ALLOWED_OPS = {
    ast.Add: op.add,
    ast.Sub: op.sub,
    ast.Mult: op.mul,
    ast.Div: op.truediv,
    ast.FloorDiv: op.floordiv,
    ast.Mod: op.mod,
    ast.Pow: op.pow,
    ast.USub: op.neg,
    ast.UAdd: op.pos,
}

def _safe_eval(node):
    if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
        return node.value
    if isinstance(node, ast.Num):
        return node.n
    if isinstance(node, ast.BinOp) and type(node.op) in _ALLOWED_OPS:
        left = _safe_eval(node.left)
        right = _safe_eval(node.right)
        return _ALLOWED_OPS[type(node.op)](left, right)
    if isinstance(node, ast.UnaryOp) and type(node.op) in _ALLOWED_OPS:
        operand = _safe_eval(node.operand)
        return _ALLOWED_OPS[type(node.op)](operand)
    raise ValueError("지원하지 않는 수식입니다.")

def try_calculate(text: str):
    raw = (text or "").strip()
    lowered = raw.lower()

    prefixes = ["계산해줘", "계산", "연산", "solve", "calc"]
    expr = raw
    for p in prefixes:
        if lowered.startswith(p):
            expr = raw[len(p):].strip()
            break

    expr = expr.replace("×", "*").replace("÷", "/").replace("^", "**")
    expr = expr.replace(" ", "")

    if not expr:
        return None

    if not re.fullmatch(r"[0-9\.\+\-\*\/%\(\)]{1,200}", expr):
        return None

    try:
        parsed = ast.parse(expr, mode="eval")
        result = _safe_eval(parsed.body)
        if isinstance(result, float) and result.is_integer():
            result = int(result)
        return f"계산 결과는 {result} 입니다."
    except Exception:
        return "수식을 계산하지 못했습니다. 예: 12*(3+4) / 2"

def parse_slide_count(text):
    m = re.search(r"(\d+)\s*장", text)
    if m:
        return int(m.group(1))
    return None

def build_ai_messages(user_id, username, latest_text):
    state = get_state(user_id)
    topic = state.get("topic")
    task = state.get("task")
    options = state.get("options", {})
    system = (
        "너는 PICK이라는 한국어 AI 작업 도우미다. "
        "아주 짧고 빠르게 답하고, 사용자의 이전 작업을 이어서 처리한다. "
        "PPT 작업 중이면 현재 상태를 고려해 수정 방향을 설명한다. "
        "예시는 전체 규칙으로 일반화하고, 주제의 정의/핵심 특징/기본 요소/활용 또는 중요성/정리를 기본 뼈대로 삼는다."
    )
    state_text = f"현재 사용자: {username}\n현재 task: {task}\n현재 topic: {topic}\n현재 options: {options}"
    return [
        {"role": "system", "content": system},
        {"role": "system", "content": state_text},
        {"role": "user", "content": latest_text},
    ]

def ai_text_reply(user_id, username, latest_text):
    if OpenAI is None:
        raise RuntimeError("openai 패키지가 설치되어 있지 않습니다.")
    client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))
    if not os.environ.get("OPENAI_API_KEY"):
        raise RuntimeError("OPENAI_API_KEY 환경변수가 없습니다.")
    messages = build_ai_messages(user_id, username, latest_text)
    response = client.responses.create(
        model=os.environ.get("OPENAI_MODEL", "gpt-5-mini"),
        input=messages
    )
    return response.output_text


def summarize_state(state: dict) -> str:
    task = state.get("task")
    topic = state.get("topic")
    options = state.get("options", {})
    if not task:
        return "현재 진행 중인 작업 없음"
    parts = [f"작업={task}"]
    if topic:
        parts.append(f"주제={topic}")
    if options:
        shown = []
        for k, v in options.items():
            if v not in [None, False, "", {}]:
                shown.append(f"{k}={v}")
        if shown:
            parts.append("옵션=" + ", ".join(shown))
    return " / ".join(parts)

def extract_topic_for_ppt(text: str) -> str:
    cleaned = re.sub(r"ppt|피피티|프레젠테이션|발표자료|만들어줘|만들기|만들어|작성해줘|생성해줘|준비해줘", " ", text, flags=re.I)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned or ""

def classify_intent_v2(text: str, state: dict) -> str:
    t = (text or "").strip()
    lowered = t.lower()
    if not t:
        return "empty"
    if try_calculate(t):
        return "calc"
    if "사이트" in t and any(x in t for x in ["추천", "찾아", "알려"]):
        return "site_recommend"
    if any(x in lowered for x in ["ppt", "피피티", "프레젠테이션"]) and any(x in t for x in ["만들", "생성", "작성", "준비"]):
        return "ppt_start"
    if state.get("task") == "ppt":
        if any(x in t for x in ["장", "이미지", "사진", "제목", "만들어", "생성", "저장", "파일", "그거", "이거", "아까", "계속"]):
            return "ppt_edit"
    if "이미지" in t and any(x in t for x in ["분석", "열어", "화면", "기능"]):
        return "image_tool"
    if ("동영상" in t or "영상" in t) and any(x in t for x in ["분석", "열어", "화면", "기능"]):
        return "video_tool"
    if "파일" in t and any(x in t for x in ["분석", "열어", "화면", "기능"]):
        return "file_tool"
    if ("exe" in lowered or "실행파일" in t) and any(x in t for x in ["만들", "열어", "기능", "도우미"]):
        return "exe_tool"
    if any(x in t for x in ["생각해서", "단계적으로", "정리해서", "검토해서", "설명", "왜", "이유", "어떻게", "비교", "차이", "요약", "정리"]):
        return "thinking_chat"
    return "chat"

def need_clarification(intent: str, text: str, state: dict):
    if intent == "ppt_start":
        topic = extract_topic_for_ppt(text)
        if not topic:
            return "PPT 주제를 먼저 말씀해 주세요."
    if intent == "ppt_edit" and state.get("task") != "ppt":
        return "현재 진행 중인 PPT 작업이 없습니다. 먼저 주제를 말씀해 주세요."
    if any(x in text for x in ["그거", "이거", "아까", "계속"]) and not state.get("task"):
        return "무엇을 이어서 처리할지 먼저 말씀해 주세요."
    return None

def make_plan_v2(intent: str, text: str, state: dict):
    if intent == "calc":
        return ["수식 추출", "안전 계산", "결과 반환"]
    if intent == "site_recommend":
        return ["주제 파악", "추천 목록 선택", "링크 반환"]
    if intent == "ppt_start":
        return ["주제 추출", "PPT 상태 저장", "장수와 이미지 여부 확인"]
    if intent == "ppt_edit":
        return ["현재 PPT 상태 확인", "옵션 수정", "생성 여부 판단"]
    if intent == "image_tool":
        return ["이미지 분석 화면 열기"]
    if intent == "video_tool":
        return ["동영상 분석 화면 열기"]
    if intent == "file_tool":
        return ["파일 분석 화면 열기"]
    if intent == "exe_tool":
        return ["EXE 도우미 화면 열기"]
    if intent == "thinking_chat":
        return ["질문 핵심 파악", "답변 구조화", "짧고 분명하게 설명"]
    return ["일반 대화 처리"]

def choose_tool(intent: str) -> str:
    mapping = {
        "calc": "calculator",
        "site_recommend": "site_recommend",
        "ppt_start": "ppt",
        "ppt_edit": "ppt",
        "image_tool": "image_ui",
        "video_tool": "video_ui",
        "file_tool": "file_ui",
        "exe_tool": "exe_ui",
        "thinking_chat": "reasoning",
        "chat": "chat",
        "empty": "chat",
    }
    return mapping.get(intent, "chat")

def build_structured_answer_v2(text: str, state: dict) -> str:
    base = simple_reasoned_reply(text)
    return (
        f"요청 이해: {text}\n"
        f"현재 상태: {summarize_state(state)}\n"
        "처리 계획:\n"
        "1. 질문의 핵심을 정리합니다.\n"
        "2. 필요한 내용을 구조화합니다.\n"
        "3. 핵심만 분명하게 설명합니다.\n\n"
        f"결과: {base}"
    )

def review_reply_v2(reply: str, state: dict) -> str:
    if not reply or not reply.strip():
        return "응답을 만들지 못했습니다."
    bad_phrases = ["이어서 진행하겠습니다", '말씀하신 "']
    for bad in bad_phrases:
        if bad in reply:
            return "질문을 이해했습니다. 조금 더 구체적으로 말씀해 주세요."
    if state.get("task") == "ppt":
        slides = state.get("options", {}).get("slides")
        if not slides and "몇 장" not in reply and "장으로" not in reply:
            return reply + "\n\n몇 장으로 만들지도 말씀해 주세요."
    return reply

def execute_plan_v2(intent: str, text: str, username: str, user_id: int):
    state = get_state(user_id)
    if intent == "empty":
        return {"reply": "말씀을 입력해 주세요."}

    if intent == "calc":
        result = try_calculate(text)
        return {"reply": result or "수식을 계산하지 못했습니다. 예: 12*(3+4)/2"}

    if intent == "site_recommend":
        sites = recommend_sites(text)
        if not sites:
            return {"reply": "해당 주제에 맞는 추천 사이트를 찾지 못했습니다."}
        reply = "추천 사이트입니다.\n\n"
        for s in sites:
            reply += f"{s['name']}\n{s['desc']}\n{s['url']}\n\n"
        return {"reply": reply}

    if intent == "ppt_start":
        topic = extract_topic_for_ppt(text)
        state["task"] = "ppt"
        state["topic"] = topic
        state["options"] = {"slides": None, "image": False, "title": topic}
        state["mode"] = "working"
        return {"reply": f'"{topic}" PPT 작업을 시작했습니다. 몇 장으로 만들까요? 이미지 포함 여부도 말씀해 주세요.'}

    if intent == "ppt_edit":
        if state.get("task") != "ppt":
            return {"reply": "현재 진행 중인 PPT 작업이 없습니다. 먼저 주제를 말씀해 주세요."}
        slide_count = extract_slide_count(text)
        if slide_count:
            state["options"]["slides"] = slide_count
            return {"reply": f"슬라이드를 {slide_count}장으로 설정했습니다."}
        if "이미지" in text or "사진" in text:
            state["options"]["image"] = True
            return {"reply": "이미지 포함으로 반영했습니다."}
        if "제목" in text and any(x in text for x in ["바꿔", "변경"]):
            title = text.replace("제목", "").replace("바꿔줘", "").replace("변경", "").strip()
            if title:
                state["options"]["title"] = title
                return {"reply": f'제목을 "{title}"로 반영했습니다.'}
            return {"reply": "제목 변경 요청을 반영했습니다."}
        if any(x in text for x in ["만들어", "생성", "저장", "파일"]):
            topic = state.get("topic") or "주제"
            options = state.get("options", {})
            if not options.get("slides"):
                return {"reply": "몇 장으로 만들지 먼저 말씀해 주세요."}
            try:
                filename = create_ppt(topic, options)
            except Exception:
                return {"reply": "PPT 생성 중 문제가 생겼습니다. 제목이나 장수를 다시 확인해 주세요."}
            state["mode"] = "done"
            return {"reply": f'"{topic}" PPT 생성을 완료했습니다.', "ppt_url": url_for("download_generated", filename=filename)}
        if any(x in text for x in ["그거", "이거", "아까", "계속"]):
            return {"reply": f'현재 "{state.get("topic")}" PPT 작업 중입니다. 장수, 제목, 이미지 여부, 생성 여부를 말씀해 주세요.'}
        return {"reply": "PPT 옵션을 수정 중입니다. 장수, 제목, 이미지 여부, 생성 여부를 말씀해 주세요."}

    if intent == "image_tool":
        return {"reply": "이미지 분석 화면을 열겠습니다.", "ui_action": "open_image"}
    if intent == "video_tool":
        return {"reply": "동영상 분석 화면을 열겠습니다.", "ui_action": "open_video"}
    if intent == "file_tool":
        return {"reply": "파일 분석 화면을 열겠습니다.", "ui_action": "open_file"}
    if intent == "exe_tool":
        return {"reply": "EXE 도우미 화면을 열겠습니다.", "ui_action": "open_exe"}

    if intent == "thinking_chat":
        return {"reply": build_structured_answer_v2(text, state)}

    return {"reply": local_chat_reply(text, username)}

def reply_v2(text: str, message_count: int, username: str, user_id: int):
    state = get_state(user_id)
    update_history(state, "user", text)
    intent = classify_intent_v2(text, state)
    state["last_intent"] = intent

    clarification = need_clarification(intent, text, state)
    if clarification:
        update_history(state, "assistant", clarification)
        return {"reply": clarification}

    state["summary"] = summarize_state(state)
    state["plan"] = make_plan_v2(intent, text, state)
    state["tool"] = choose_tool(intent)

    result = execute_plan_v2(intent, text, username, user_id)
    result["reply"] = review_reply_v2(result.get("reply", ""), state)

    state["summary"] = summarize_state(state)
    update_history(state, "assistant", result["reply"])
    return result




def reply(text, message_count, username, user_id):
    normalized_text, corrections = pick_typo.normalize(text)

    if corrections and corrections[0].get("type") == "learned":
        reply_text = f"{corrections[0]['wrong']} → {corrections[0]['correct']}로 기억했습니다. 앞으로 같은 표현은 자동으로 보정하겠습니다."
        return {"reply": final_review(reply_text, text, "오타 학습", "", False)}

    tools = {
        "create_ppt": create_ppt,
        "llm": pick_llm,
    }

    result = pick_engine.handle(normalized_text, username, user_id, tools)

    state = pick_engine.get_state(user_id) if hasattr(pick_engine, "get_state") else {}
    intent = state.get("last_intent", "")
    state_summary = state.get("summary", "")

    thinking = any(k in text for k in ["생각", "단계적으로", "검토", "정리해서", "이유", "왜"])

    if corrections:
        note = "입력 보정: " + ", ".join([f"{c['wrong']}→{c['correct']}" for c in corrections])
        result["reply"] = f"{note}\\n\\n{result.get('reply', '')}"

    result["reply"] = final_review(
        result.get("reply", ""),
        user_text=normalized_text,
        intent=intent,
        state_summary=state_summary,
        thinking=thinking
    )

    result["reply"] = result.get("reply", "").replace("Plcker", "PICK").replace("plcker", "PICK").replace("Picker", "PICK")

    result["reply"] = enforce_polite_pick_reply(result.get("reply", ""))

    if result.get("ppt_filename"):
        result["ppt_url"] = url_for("download_generated", filename=result["ppt_filename"])

    return result


def enforce_polite_pick_reply(text):
    t = str(text or "").strip()
    if not t:
        return "죄송합니다. 응답을 만들지 못했습니다."

    fixes = {
        "뭘 도와줄까?": "무엇을 도와드릴까요?",
        "도와줄게": "도와드리겠습니다",
        "말해줘": "말씀해 주세요",
        "알려줘": "알려주세요",
        "해줘": "해 주세요",
        "응": "네",
        "ㅇ": "네",
    }
    for a, b in fixes.items():
        t = t.replace(a, b)

    t = t.replace("Plcker", "PICK").replace("Picker", "PICK").replace("plcker", "PICK")

    # 같은 줄 반복 제거
    lines = []
    seen = set()
    for line in t.splitlines():
        key = line.strip()
        if not key:
            lines.append(line)
            continue
        if key in seen:
            continue
        seen.add(key)
        lines.append(line)
    t = "\n".join(lines).strip()

    return t


@app.route("/")
@login_required
def index():
    if not logged_in():
        return render_template("auth.html")
    return render_template("app.html", username=session["username"], is_admin=(session.get("username") == ADMIN_USERNAME))

@app.get("/api/bootstrap")
def api_bootstrap():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    user_id = current_user_id()
    conn = db()
    projects = [dict(r) for r in conn.execute("SELECT * FROM projects WHERE user_id=? ORDER BY id DESC", (user_id,)).fetchall()]
    requests_rows = [dict(r) for r in conn.execute("SELECT * FROM requests WHERE user_id=? ORDER BY id DESC", (user_id,)).fetchall()]
    conn.close()
    return jsonify({
        "ok": True,
        "username": session["username"],
        "chats": get_user_chats(user_id),
        "projects": projects,
        "requests": requests_rows,
    })

def user_owns_chat(chat_id: int, user_id: int) -> bool:
    conn = db()
    row = conn.execute("SELECT id FROM chats WHERE id=? AND user_id=?", (chat_id, user_id)).fetchone()
    conn.close()
    return row is not None


@app.post("/api/chat/new")
def api_chat_new():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    chat_id = create_chat(current_user_id(), with_greeting=True)
    return jsonify({"ok": True, "chat_id": chat_id, "messages": get_messages(chat_id), "chats": get_user_chats(current_user_id())})

@app.get("/api/chat/<int:chat_id>")
def api_chat_get(chat_id):
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    if not user_owns_chat(chat_id, current_user_id()):
        return jsonify({"ok": False, "error": "채팅을 찾을 수 없습니다."}), 404
    return jsonify({"ok": True, "messages": get_messages(chat_id)})

@app.post("/api/chat/<int:chat_id>/delete")
def api_chat_delete(chat_id):
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    if not user_owns_chat(chat_id, current_user_id()):
        return jsonify({"ok": False, "error": "채팅을 찾을 수 없습니다."}), 404
    conn = db()
    conn.execute("DELETE FROM messages WHERE chat_id=?", (chat_id,))
    conn.execute("DELETE FROM chats WHERE id=? AND user_id=?", (chat_id, current_user_id()))
    conn.commit()
    conn.close()
    return jsonify({"ok": True, "chats": get_user_chats(current_user_id())})

@app.post("/api/chat/<int:chat_id>/send")
def api_chat_send(chat_id):
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    if not user_owns_chat(chat_id, current_user_id()):
        return jsonify({"ok": False, "error": "채팅을 찾을 수 없습니다."}), 404
    text = request.form.get("message", "").strip()
    if not text:
        return jsonify({"ok": False, "error": "빈 메시지입니다."}), 400
    conn = db()
    conn.execute("INSERT INTO messages(chat_id, role, content, created_at) VALUES(?,?,?,?)", (chat_id, "user", text, now()))
    conn.commit()
    update_chat_title(chat_id)
    messages = get_messages(chat_id)
    result = reply(text, len(messages), session["username"], current_user_id())
    conn = db()
    conn.execute("INSERT INTO messages(chat_id, role, content, created_at) VALUES(?,?,?,?)", (chat_id, "bot", result["reply"], now()))
    conn.commit()
    conn.close()
    out = {"ok": True, "messages": get_messages(chat_id), "chats": get_user_chats(current_user_id())}
    out.update(result)
    return jsonify(out)

@app.post("/api/project")
def api_project():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    title = request.form.get("title", "").strip()
    description = request.form.get("description", "").strip()
    conn = db()
    conn.execute("INSERT INTO projects(user_id, title, description, created_at) VALUES(?,?,?,?)",
                 (current_user_id(), title, description, now()))
    conn.commit()
    rows = [dict(r) for r in conn.execute("SELECT * FROM projects WHERE user_id=? ORDER BY id DESC", (current_user_id(),)).fetchall()]
    conn.close()
    return jsonify({"ok": True, "projects": rows})

@app.post("/api/request")
def api_request():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    content = request.form.get("content", "").strip()
    conn = db()
    conn.execute("INSERT INTO requests(user_id, content, created_at) VALUES(?,?,?)",
                 (current_user_id(), content, now()))
    conn.commit()
    rows = [dict(r) for r in conn.execute("SELECT * FROM requests WHERE user_id=? ORDER BY id DESC", (current_user_id(),)).fetchall()]
    conn.close()
    return jsonify({"ok": True, "requests": rows})

@app.post("/api/analyze/image")
def api_analyze_image():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    f = request.files.get("file")
    if not f:
        return jsonify({"ok": False, "error": "파일이 없습니다."}), 400
    path = UPLOADS / secure_filename(f.filename)
    f.save(path)
    return jsonify({"ok": True, "result": analyze_image(path)})

@app.post("/api/analyze/video")
def api_analyze_video():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    f = request.files.get("file")
    if not f:
        return jsonify({"ok": False, "error": "파일이 없습니다."}), 400
    path = UPLOADS / secure_filename(f.filename)
    f.save(path)
    result = analyze_file(path)
    result["note"] = "현재 버전은 동영상 메타데이터 중심 분석입니다. 의미 요약형 영상 분석은 별도 엔진이 필요합니다."
    return jsonify({"ok": True, "result": result})

@app.post("/api/analyze/file")
def api_analyze_file():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    f = request.files.get("file")
    if not f:
        return jsonify({"ok": False, "error": "파일이 없습니다."}), 400
    path = UPLOADS / secure_filename(f.filename)
    f.save(path)
    return jsonify({"ok": True, "result": analyze_file(path)})

@app.post("/api/exe/main.py")
def api_exe_main():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    code = request.form.get("code", 'print("Hello from PICK")')
    return send_file(io.BytesIO(code.encode("utf-8")), mimetype="text/x-python", as_attachment=True, download_name="main.py")

@app.post("/api/exe/build.bat")
def api_exe_bat():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    name = re.sub(r"[^\w\-]", "_", request.form.get("name", "app")) or "app"
    content = (
        "@echo off\r\n"
        "py -m venv .venv\r\n"
        "call .venv\\Scripts\\activate\r\n"
        "pip install pyinstaller\r\n"
        f"pyinstaller --onefile --name {name} main.py\r\n"
        "pause\r\n"
    )
    return send_file(io.BytesIO(content.encode("utf-8")), mimetype="text/plain", as_attachment=True, download_name="build.bat")

@app.post("/api/exe/readme")
def api_exe_readme():
    auth_err = require_login_json()
    if auth_err:
        return auth_err
    name = re.sub(r"[^\w\-]", "_", request.form.get("name", "app")) or "app"
    content = (
        f"# {name} EXE 빌드 안내\n\n"
        "1. main.py 와 build.bat 를 같은 폴더에 저장\n"
        "2. Windows에서 build.bat 실행\n"
        f"3. dist 폴더 안에 {name}.exe 생성\n"
    )
    return send_file(io.BytesIO(content.encode("utf-8")), mimetype="text/plain", as_attachment=True, download_name="README_EXE.txt")

@app.get("/generated/<path:filename>")
def download_generated(filename):
    return send_from_directory(GENERATED, filename, as_attachment=True)


@app.route("/api/paste/analyze", methods=["POST"])
def api_paste_analyze():
    if "file" not in request.files:
        return jsonify({"ok": False, "error": "파일이 없습니다."})

    file = request.files["file"]
    filename = secure_filename(file.filename or "pasted_file")
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    data = file.read()
    size = len(data)

    image_exts = {"png", "jpg", "jpeg", "gif", "webp", "bmp"}
    video_exts = {"mp4", "mov", "avi", "mkv", "webm"}
    doc_exts = {"txt", "md", "csv", "json", "pdf", "docx", "pptx", "xlsx", "py", "js", "html", "css"}

    if ext in image_exts:
        kind = "image"
        message = (
            "이미지 파일을 붙여넣기로 받았습니다.\\n"
            f"파일명: {filename}\\n"
            f"크기: {size:,} bytes\\n"
            "현재는 기본 이미지 파일 정보 분석을 제공합니다. "
            "LLaVA 같은 비전 모델을 연결하면 실제 이미지 내용 설명까지 가능합니다."
        )
    elif ext in video_exts:
        kind = "video"
        message = (
            "동영상 파일을 붙여넣기로 받았습니다.\\n"
            f"파일명: {filename}\\n"
            f"크기: {size:,} bytes\\n"
            "현재는 동영상 파일 구조 정보를 분석합니다. "
            "프레임 분석 모델을 연결하면 장면 설명까지 확장할 수 있습니다."
        )
    elif ext in doc_exts:
        kind = "file"
        sample = ""
        if ext in {"txt", "md", "csv", "json", "py", "js", "html", "css"}:
            try:
                sample = data[:1500].decode("utf-8", errors="ignore")
            except Exception:
                sample = ""
        message = (
            "파일을 붙여넣기로 받았습니다.\\n"
            f"파일명: {filename}\\n"
            f"확장자: {ext or '없음'}\\n"
            f"크기: {size:,} bytes"
        )
        if sample:
            message += "\\n\\n미리보기:\\n" + sample
    else:
        kind = "file"
        message = (
            "붙여넣은 파일을 받았습니다.\\n"
            f"파일명: {filename}\\n"
            f"확장자: {ext or '없음'}\\n"
            f"크기: {size:,} bytes\\n"
            "지원하지 않는 형식일 수 있지만 기본 파일 정보는 확인했습니다."
        )

    return jsonify({
        "ok": True,
        "kind": kind,
        "filename": filename,
        "size": size,
        "result": message
    })


def save_service_message(conversation_id, role, content):
    if not session.get("user_id"):
        return
    try:
        conn = db()
        owner = conn.execute(
            "SELECT id FROM conversations WHERE id=? AND user_id=?",
            (conversation_id, session["user_id"])
        ).fetchone()
        if not owner:
            conn.close()
            return
        conn.execute(
            "INSERT INTO messages(conversation_id, role, content, created_at) VALUES (?, ?, ?, ?)",
            (conversation_id, role, content, now_iso())
        )
        conn.execute("UPDATE conversations SET updated_at=? WHERE id=?", (now_iso(), conversation_id))
        conn.commit()
        conn.close()
    except Exception as e:
        service_log("ERROR", f"save message failed: {e}")

# ===== PICK SERVICE ROUTES =====


# ===== PICK SERVICE CLEAN AUTH ROUTES =====
@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "GET":
        return render_template("register.html")

    username = request.form.get("username", "").strip()
    password = request.form.get("password", "").strip()
    password2 = request.form.get("password2", password).strip()

    if has_korean(username) or has_korean(password) or has_korean(password2):
        return render_template("register.html", error="아이디 또는 비밀번호에 한글을 사용할 수 없습니다.")

    if password != password2:
        return render_template("register.html", error="비밀번호 확인이 일치하지 않습니다.")

    if not valid_account_text(username, allow_symbols=False):
        return render_template("register.html", error="아이디는 영어, 숫자, _, - 만 사용할 수 있으며 2~32자여야 합니다.")

    if not valid_account_text(password, allow_symbols=True):
        return render_template("register.html", error="비밀번호는 영어, 숫자, 일부 특수문자만 사용할 수 있으며 4~64자여야 합니다.")

    conn = db()
    try:
        conn.execute(
            "INSERT INTO users(username, password_hash, created_at) VALUES (?, ?, ?)",
            (username, generate_password_hash(password), now_iso())
        )
        conn.commit()
    except sqlite3.IntegrityError:
        conn.close()
        return render_template("register.html", error="이미 사용 중인 아이디입니다.")

    user = conn.execute("SELECT id FROM users WHERE username=?", (username,)).fetchone()
    conn.close()

    session["user_id"] = user["id"]
    session["username"] = username
    service_log("INFO", f"new user registered: {username}")
    return redirect(url_for("index"))

@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "GET":
        return render_template("login.html")

    username = request.form.get("username", "").strip()
    password = request.form.get("password", "").strip()

    if has_korean(username) or has_korean(password):
        return render_template("login.html", error="아이디 또는 비밀번호에 한글을 사용할 수 없습니다.")

    conn = db()
    user = conn.execute("SELECT * FROM users WHERE username=?", (username,)).fetchone()
    conn.close()

    if not user or not check_password_hash(user["password_hash"], password):
        return render_template("login.html", error="아이디 또는 비밀번호가 올바르지 않습니다.")

    session["user_id"] = user["id"]
    session["username"] = user["username"]
    service_log("INFO", f"user login: {username}")
    return redirect(url_for("index"))

@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))

@app.route("/admin/status")
@login_required
def admin_status():
    if session.get("username") != ADMIN_USERNAME:
        return "관리자만 접근할 수 있습니다.", 403
    conn = db()
    users = conn.execute("SELECT COUNT(*) AS c FROM users").fetchone()["c"]
    convs = conn.execute("SELECT COUNT(*) AS c FROM conversations").fetchone()["c"]
    msgs = conn.execute("SELECT COUNT(*) AS c FROM messages").fetchone()["c"]
    logs = conn.execute("SELECT level, message, created_at FROM service_logs ORDER BY id DESC LIMIT 30").fetchall()
    conn.close()
    return render_template("admin_status.html", users=users, convs=convs, msgs=msgs, logs=logs)

# Gunicorn/Render로 실행할 때도 데이터베이스와 관리자 계정을 준비합니다.
init_db()

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=False)
