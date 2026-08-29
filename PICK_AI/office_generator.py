from __future__ import annotations

import json
import re
import uuid
from pathlib import Path
from typing import Any

from config import GENERATED_DIR
from pick_llm import PickLLMRouter

try:
    from docx import Document
except Exception:
    Document = None

try:
    from openpyxl import Workbook
except Exception:
    Workbook = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

CREATE_WORDS = ("만들", "생성", "작성", "제작", "파일로", "정리해줘", "정리해 주세요")


def detect_office_kind(text: str) -> str | None:
    raw = str(text or "")
    lower = raw.lower()
    if not any(word in lower for word in CREATE_WORDS):
        return None
    if any(x in lower for x in ("ppt", "pptx", "파워포인트", "프레젠테이션", "프리젠테이션", "프라이전체이션", "프라이젠테이션", "발표자료", "발표 자료")):
        return "pptx"
    if any(x in lower for x in ("xlsx", "엑셀", "excel", "스프레드시트", "표 파일")):
        return "xlsx"
    if any(x in lower for x in ("docx", "워드", "word", "워드 파일", "문서 파일")):
        return "docx"
    return None


def _safe_title(text: str, fallback: str) -> str:
    clean = re.sub(
        r"(pptx?|파워포인트|프레젠테이션|프리젠테이션|프라이전체이션|프라이젠테이션|발표\s*자료|엑셀|excel|xlsx|"
        r"워드|word|docx|파일|만들어\s*줘|만들어줘|만들어\s*주세요|생성|작성|제작)",
        " ", str(text or ""), flags=re.I,
    )
    clean = re.sub(r"\s+", " ", clean).strip(" ?!.")
    return clean[:80] or fallback


def _llm(prompt: str) -> str:
    return PickLLMRouter().generate(prompt)


def _extract_json(text: str) -> dict[str, Any] | None:
    value = str(text or "").strip()
    value = re.sub(r"^```(?:json)?\s*", "", value, flags=re.I)
    value = re.sub(r"\s*```$", "", value)
    try:
        data = json.loads(value)
        return data if isinstance(data, dict) else None
    except Exception:
        pass
    start, end = value.find("{"), value.rfind("}")
    if start >= 0 and end > start:
        try:
            data = json.loads(value[start:end + 1])
            return data if isinstance(data, dict) else None
        except Exception:
            return None
    return None


def _pptx(text: str) -> Path:
    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다.")
    title = _safe_title(text, "PICK 프레젠테이션")
    prompt = f"""다음 요청으로 발표자료 내용을 만드세요.
반드시 JSON만 출력하세요.
형식: {{"title":"제목","subtitle":"부제","slides":[{{"title":"슬라이드 제목","bullets":["핵심 1","핵심 2"]}}]}}
슬라이드는 5~10장, 각 슬라이드는 핵심 bullet 2~5개로 작성하세요.
요청: {text}
"""
    data = _extract_json(_llm(prompt)) or {}
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = str(data.get("title") or title)
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = str(data.get("subtitle") or "PICK이 생성한 발표자료")

    slides = data.get("slides")
    if not isinstance(slides, list) or not slides:
        slides = [
            {"title": "요청", "bullets": [text]},
            {"title": "핵심 내용", "bullets": ["요청 내용을 기준으로 발표 구조를 만들었습니다."]},
            {"title": "다음 단계", "bullets": ["필요한 수치, 이미지, 출처를 추가해 완성도를 높이세요."]},
        ]
    for item in slides[:12]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str((item or {}).get("title") or "내용")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = (item or {}).get("bullets") or []
        for index, bullet in enumerate(bullets[:6]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0

    path = GENERATED_DIR / f"pick_presentation_{uuid.uuid4().hex[:10]}.pptx"
    prs.save(path)
    return path


def _docx(text: str) -> Path:
    if Document is None:
        raise RuntimeError("python-docx가 설치되어 있지 않습니다.")
    title = _safe_title(text, "PICK 문서")
    prompt = f"""다음 요청을 바탕으로 전문적인 문서 내용을 작성하세요.
Markdown 형식으로 제목, 소제목, 본문을 작성하고 불필요한 서론은 줄이세요.
요청: {text}
"""
    content = _llm(prompt)
    doc = Document()
    doc.add_heading(title, level=0)
    for line in str(content or text).splitlines():
        line = line.rstrip()
        if not line:
            continue
        if line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif re.match(r"^[-*]\s+", line):
            doc.add_paragraph(re.sub(r"^[-*]\s+", "", line), style="List Bullet")
        else:
            doc.add_paragraph(line)
    path = GENERATED_DIR / f"pick_document_{uuid.uuid4().hex[:10]}.docx"
    doc.save(path)
    return path


def _xlsx(text: str) -> Path:
    if Workbook is None:
        raise RuntimeError("openpyxl이 설치되어 있지 않습니다.")
    prompt = f"""다음 요청을 엑셀 표로 정리하세요.
반드시 JSON만 출력하세요.
형식: {{"sheet":"시트명","headers":["열1","열2"],"rows":[["값1","값2"]]}}
열은 2~8개, 행은 요청에 맞게 작성하세요. 확인되지 않은 숫자는 만들지 마세요.
요청: {text}
"""
    data = _extract_json(_llm(prompt)) or {}
    wb = Workbook()
    ws = wb.active
    ws.title = str(data.get("sheet") or "PICK")[:31]
    headers = data.get("headers")
    rows = data.get("rows")
    if not isinstance(headers, list) or not headers:
        headers = ["항목", "내용"]
    if not isinstance(rows, list) or not rows:
        rows = [["요청", text], ["설명", "PICK이 생성한 기본 표입니다."]]
    ws.append([str(x) for x in headers])
    for row in rows[:500]:
        if isinstance(row, list):
            ws.append([x if isinstance(x, (int, float)) else str(x) for x in row[:len(headers)]])
    for column in ws.columns:
        max_len = max(len(str(cell.value or "")) for cell in column)
        ws.column_dimensions[column[0].column_letter].width = min(max(max_len + 2, 10), 50)
    path = GENERATED_DIR / f"pick_spreadsheet_{uuid.uuid4().hex[:10]}.xlsx"
    wb.save(path)
    return path


def create_office_file(text: str, kind: str | None = None) -> dict[str, str]:
    kind = kind or detect_office_kind(text)
    if kind == "pptx":
        path = _pptx(text); label = "PowerPoint"
    elif kind == "xlsx":
        path = _xlsx(text); label = "Excel"
    elif kind == "docx":
        path = _docx(text); label = "Word"
    else:
        raise ValueError("지원하는 문서 형식을 찾지 못했습니다.")
    return {"kind": kind, "label": label, "filename": path.name, "url": f"/generated/{path.name}"}


def smoke_test_files() -> list[str]:
    created = []
    if Presentation is not None:
        prs = Presentation(); prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = "PICK"
        p = GENERATED_DIR / "_pick_smoke.pptx"; prs.save(p); created.append(str(p))
    if Document is not None:
        doc = Document(); doc.add_paragraph("PICK")
        p = GENERATED_DIR / "_pick_smoke.docx"; doc.save(p); created.append(str(p))
    if Workbook is not None:
        wb = Workbook(); wb.active["A1"] = "PICK"
        p = GENERATED_DIR / "_pick_smoke.xlsx"; wb.save(p); created.append(str(p))
    for value in created:
        Path(value).unlink(missing_ok=True)
    return created
